from pathlib import Path
from typing import List, Tuple
from pptx import Presentation as PPTXPresentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from PIL import Image, ImageDraw, ImageFont
import io
import os

from ppt2vid.models.slide import Slide, Presentation as PresentationModel

class PPTConverter:
    """Handles conversion of PowerPoint presentations to images."""

    def __init__(self, output_dir: Path):
        """Initialize the converter with output directory."""
        self.output_dir = Path(output_dir)
        self.output_dir.mkdir(parents=True, exist_ok=True)
        
        # Try to load a font that supports UTF-8
        try:
            self.font = ImageFont.truetype("/System/Library/Fonts/Arial Unicode.ttf", 20)
        except OSError:
            try:
                self.font = ImageFont.truetype("/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf", 20)
            except OSError:
                self.font = ImageFont.load_default()

    def convert_to_images(self, ppt_path: Path) -> PresentationModel:
        """Convert PowerPoint presentation to images.
        
        Args:
            ppt_path: Path to the PowerPoint file
            
        Returns:
            Presentation object containing slide information
        """
        prs = PPTXPresentation(ppt_path)
        slides: List[Slide] = []
        
        for idx, slide in enumerate(prs.slides):
            # Create slide-specific output directory
            slide_dir = self.output_dir / f"slide_{idx + 1}"
            slide_dir.mkdir(exist_ok=True)
            
            # Save slide as image
            image_path = slide_dir / "slide.png"
            self._save_slide_as_image(prs, slide, image_path)
            
            # Create slide object
            slides.append(Slide(
                slide_id=idx + 1,
                image_path=image_path
            ))
        
        return PresentationModel(
            slides=slides,
            title=ppt_path.stem,
            output_dir=self.output_dir
        )

    def _get_text_from_shape(self, shape) -> str:
        """Extract text from a shape, handling potential encoding issues.
        
        Args:
            shape: PowerPoint shape object
            
        Returns:
            Text content of the shape
        """
        try:
            if not hasattr(shape, 'text_frame'):
                return ""
            
            text_parts = []
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    # Encode and decode using UTF-8 to handle Unicode characters
                    text = run.text.encode('utf-8', errors='ignore').decode('utf-8')
                    text_parts.append(text)
            return "\n".join(text_parts)
        except Exception as e:
            print(f"Error processing shape: {e}")
            return ""

    def _get_shape_color(self, shape) -> Tuple[int, int, int]:
        """Get the color of a shape.
        
        Args:
            shape: PowerPoint shape object
            
        Returns:
            RGB color tuple
        """
        try:
            if hasattr(shape, 'fill'):
                if shape.fill.type:
                    if hasattr(shape.fill.fore_color, 'rgb'):
                        rgb = shape.fill.fore_color.rgb
                        return (rgb >> 16) & 255, (rgb >> 8) & 255, rgb & 255
        except Exception:
            pass
        return 0, 0, 0  # Default to black

    def _save_slide_as_image(self, prs, slide, output_path: Path) -> None:
        """Save a single slide as an image.
        
        Args:
            prs: PowerPoint presentation object
            slide: PowerPoint slide object
            output_path: Path where the image should be saved
        """
        # Standard PowerPoint slide dimensions (in pixels at 96 dpi)
        # 16:9 aspect ratio at 1920x1080
        width = 1920
        height = 1080
        
        # Create a new image with white background
        image = Image.new('RGB', (width, height), 'white')
        draw = ImageDraw.Draw(image)
        
        # Process shapes in z-order (back to front)
        for shape in slide.shapes:
            try:
                # Get shape dimensions
                left = shape.left * width / prs.slide_width
                top = shape.top * height / prs.slide_height
                right = (shape.left + shape.width) * width / prs.slide_width
                bottom = (shape.top + shape.height) * height / prs.slide_height
                
                # Handle different shape types
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    # Extract image data
                    image_stream = io.BytesIO(shape.image.blob)
                    shape_image = Image.open(image_stream)
                    
                    # Resize image to fit shape dimensions
                    shape_image = shape_image.resize((int(right - left), int(bottom - top)))
                    
                    # Paste image onto slide
                    image.paste(shape_image, (int(left), int(top)))
                else:
                    # Draw shape outline and fill
                    color = self._get_shape_color(shape)
                    draw.rectangle([left, top, right, bottom], outline='black', fill=color)
                
                # Draw text if present
                text = self._get_text_from_shape(shape)
                if text:
                    # Draw text with word wrapping
                    text_width = right - left - 10
                    y = top + 5
                    for line in text.split('\n'):
                        words = line.split()
                        line_parts = []
                        current_line = []
                        
                        for word in words:
                            current_line.append(word)
                            test_line = ' '.join(current_line)
                            bbox = draw.textbbox((0, 0), test_line, font=self.font)
                            if bbox[2] > text_width:
                                if current_line:
                                    line_parts.append(' '.join(current_line[:-1]))
                                    current_line = [word]
                        
                        if current_line:
                            line_parts.append(' '.join(current_line))
                        
                        for part in line_parts:
                            draw.text((left + 5, y), part, fill='black', font=self.font)
                            y += 25  # Increased line spacing for better readability
                        y += 5  # Paragraph spacing
            except Exception as e:
                print(f"Error processing shape: {e}")
                continue
        
        # Save the image
        image.save(output_path, 'PNG') 